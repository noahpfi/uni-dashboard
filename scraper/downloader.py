"""Step 5: download newly-added/changed priority files + extract text"""
import argparse
import io
import json
import re
import time
from pathlib import Path

from playwright.sync_api import sync_playwright
from auth import make_context, ensure_logged_in, BASE

DATA_DIR = Path(__file__).parent.parent / 'data'
EXTRACT_TYPES = {'pdf', 'pptx', 'docx'}
MAX_MB = 10


def _parse_mb(size_str: str) -> float:
    s = (size_str or '').strip().upper()
    m = re.match(r'([\d.]+)\s*(KB|MB|GB|B)?', s)
    if not m:
        return 0.0
    val, unit = float(m.group(1)), (m.group(2) or 'B')
    return val * {'GB': 1024, 'MB': 1, 'KB': 1 / 1024, 'B': 1 / (1024 ** 2)}[unit]


def _file_key(f: dict) -> str:
    return f"{f.get('folder_nid', '')}/{f.get('folder_path', '')}/{f['filename']}"


def _text_key(f: dict) -> str:
    raw = f"{f['folder_nid']}_{f['folder_path']}_{f['filename']}"
    return re.sub(r'[^\w.\-]', '_', raw) + '.txt'


def _folder_url(cid: str, nid: str, enc_path: str) -> str:
    if enc_path:
        return f"{BASE}/auth/RepositoryEntry/{cid}/CourseNode/{nid}/path%3D~~{enc_path}/0"
    return f"{BASE}/url/RepositoryEntry/{cid}/CourseNode/{nid}"


def _get_download_url(page, cid: str, f: dict) -> str | None:
    nid = f['folder_nid']
    enc_path = f['folder_path'].replace('/', '%2F')
    filename = f['filename']

    page.goto(_folder_url(cid, nid, enc_path), wait_until='domcontentloaded', timeout=60000)
    # Remove any lingering lightbox so the next wait_for_selector sees only a fresh one
    page.evaluate("() => { document.querySelectorAll('.basicLightbox').forEach(el => el.remove()); }")

    # Find row by data-drag-file attribute (avoids CSS selector escaping issues)
    target_link = None
    for row in page.query_selector_all('tr.o_folder_row'):
        if row.get_attribute('data-drag-file') == filename:
            target_link = row.query_selector('td:nth-child(2) a')
            break

    if not target_link:
        print(f"    [miss] row/link not found: {filename}")
        return None

    # Click link; force=True bypasses pointer-event interception from any residual overlay
    target_link.click(force=True)

    # Wait for the new PDF.js lightbox iframe and read its src
    try:
        iframe = page.wait_for_selector('.basicLightbox--visible iframe[src]', timeout=12000)
        src = iframe.get_attribute('src') or ''
        m = re.search(r'file=(/m/[a-f0-9]{32})', src)
        if m:
            return f"{BASE}{m.group(1)}"
        m = re.search(r'/m/([a-f0-9]{32})', src)
        if m:
            return f"{BASE}/m/{m.group(1)}"
    except Exception:
        pass

    print(f"    [warn] no lightbox iframe found: {filename}")
    return None


def _download_bytes(page, url: str) -> bytes | None:
    try:
        resp = page.request.get(url)
        if resp.ok:
            return resp.body()
        print(f"    [err] HTTP {resp.status}: {url}")
    except Exception as e:
        print(f"    [err] download: {e}")
    return None


def _extract_text(data: bytes, ftype: str) -> str:
    try:
        if ftype == 'pdf':
            import pdfplumber
            with pdfplumber.open(io.BytesIO(data)) as pdf:
                return '\n'.join(p.extract_text() or '' for p in pdf.pages).strip()
        if ftype == 'pptx':
            from pptx import Presentation
            prs = Presentation(io.BytesIO(data))
            parts = [
                shape.text_frame.text
                for slide in prs.slides
                for shape in slide.shapes
                if shape.has_text_frame
            ]
            return '\n'.join(parts).strip()
        if ftype == 'docx':
            from docx import Document
            doc = Document(io.BytesIO(data))
            return '\n'.join(p.text for p in doc.paragraphs).strip()
    except Exception as e:
        print(f"    [err] extract ({ftype}): {e}")
    return ''


def process_course(page, mat: dict, force_all: bool) -> int:
    cid = mat['course_id']
    files = mat.get('files', [])
    diff = mat.get('diff', {})

    texts_dir = DATA_DIR / 'courses' / cid / 'texts'
    texts_dir.mkdir(parents=True, exist_ok=True)

    if force_all:
        candidates = [f for f in files
                      if f.get('priority') and f.get('type', '').lower() in EXTRACT_TYPES]
    else:
        diff_keys = {_file_key(f) for f in diff.get('added', []) + diff.get('changed', [])}
        # Also include priority files not yet extracted (backfill pre-existing files)
        candidates = [f for f in files
                      if f.get('priority')
                      and f.get('type', '').lower() in EXTRACT_TYPES
                      and (_file_key(f) in diff_keys
                           or not (texts_dir / _text_key(f)).exists())]

    name = mat.get('course_name', cid)
    if candidates:
        print(f"\n[{cid}] {name[:50]} — {len(candidates)} candidate(s)")

    processed = 0
    for f in candidates:
        filename = f['filename']
        ftype = f.get('type', '').lower()
        size_mb = _parse_mb(f.get('size', ''))

        if size_mb > MAX_MB:
            print(f"    [skip] {filename} ({size_mb:.1f} MB > {MAX_MB})")
            continue

        tpath = texts_dir / _text_key(f)
        if tpath.exists():
            print(f"    [cached] {filename}")
            f['text_file'] = str(tpath.relative_to(DATA_DIR))
            processed += 1
            continue

        print(f"    [fetch] {filename} ({ftype}, {f.get('size', '?')})")
        dl_url = _get_download_url(page, cid, f)
        if not dl_url:
            continue

        data = _download_bytes(page, dl_url)
        if not data:
            continue

        text = _extract_text(data, ftype)
        if text:
            tpath.write_text(text, encoding='utf-8')
            f['text_file'] = str(tpath.relative_to(DATA_DIR))
            print(f"      → {len(text):,} chars")
        else:
            print(f"      → empty extract")

        processed += 1

    if candidates:
        print(f"  → {processed}/{len(candidates)} processed")

    # Restore text_file for files already on disk but wiped by materials.py rebuild
    restored = 0
    for f in files:
        if f.get('text_file') or f.get('type', '').lower() not in EXTRACT_TYPES:
            continue
        tpath = texts_dir / _text_key(f)
        if tpath.exists():
            f['text_file'] = str(tpath.relative_to(DATA_DIR))
            restored += 1

    return processed + restored


def run_all(page, force_all: bool = False, only_course: str = '') -> None:
    import glob
    paths = sorted(glob.glob(str(DATA_DIR / 'courses' / '*' / 'materials.json')))

    for mat_path_str in paths:
        mat_path = Path(mat_path_str)
        with open(mat_path) as fh:
            mat = json.load(fh)

        if only_course and mat['course_id'] != only_course:
            continue

        n = process_course(page, mat, force_all)
        if n > 0:
            with open(mat_path, 'w') as fh:
                json.dump(mat, fh, indent=2, ensure_ascii=False)


def main():
    ap = argparse.ArgumentParser(description='Step 5: download + extract priority files')
    ap.add_argument('--force-all', action='store_true',
                    help='Process all priority files, not just diff')
    ap.add_argument('--course', metavar='CID', default='',
                    help='Limit to one course ID')
    args = ap.parse_args()

    with sync_playwright() as p:
        browser = p.chromium.launch(headless=True)
        ctx = make_context(browser)
        page = ctx.new_page()
        ensure_logged_in(ctx, page)
        run_all(page, force_all=args.force_all, only_course=args.course)
        browser.close()

    print('\nDone.')


if __name__ == '__main__':
    main()
